"""Populate the COUNT/REACH .pptx template by replacing {{TOKEN}} placeholders.

This never redesigns, resizes, or restyles anything — it only swaps text inside
existing runs, so every color, position, and font from the template is preserved
exactly. This is the fix for chat-based tools (like Copilot) redesigning the
slide instead of populating it.
"""

import copy
from pptx import Presentation

TEMPLATE_PATH = "assets/count_template.pptx"


def _score_display(score):
    """Turn a 1-5 score (or None) into the 'X/5' string the template expects."""
    if score is None:
        return "N/A"
    return f"{score}/5"


def build_token_map(summary):
    """Convert the AI's JSON summary into a flat {token: replacement text} map."""
    reach = summary.get("reach", {})

    def field(key, subkey):
        entry = reach.get(key, {}) or {}
        return entry.get(subkey)

    total = summary.get("total_score", "N/A")

    return {
        "{{REVIEW_DATE}}": str(summary.get("review_date", "N/A")),
        "{{SUBHEADER}}": str(summary.get("subheader", "")),
        "{{ACCOUNT_SNAPSHOT}}": str(summary.get("account_snapshot", "")),
        "{{TOTAL_SCORE}}": str(total),
        "{{CLASSIFICATION}}": str(summary.get("classification", "N/A")),

        "{{R_SCORE}}": _score_display(field("return_on_assets", "score")),
        "{{R_EVIDENCE}}": str(field("return_on_assets", "evidence") or "Not Available"),

        "{{E_SCORE}}": _score_display(field("si_value_to_customer", "score")),
        "{{E_EVIDENCE}}": str(field("si_value_to_customer", "evidence") or "Not Available"),

        "{{A_SCORE}}": _score_display(field("ar_contracts", "score")),
        "{{A_EVIDENCE}}": str(field("ar_contracts", "evidence") or "Not Available"),

        "{{C_SCORE}}": _score_display(field("customer_stage_crisis", "score")),
        "{{C_EVIDENCE}}": str(field("customer_stage_crisis", "evidence") or "Not Available"),

        "{{H_SCORE}}": _score_display(field("high_performers", "score")),
        "{{H_EVIDENCE}}": str(field("high_performers", "evidence") or "Not Available"),
    }


def fill_template(summary, output_path, template_path=TEMPLATE_PATH):
    """Open the template, replace every token in-place, save to output_path."""
    token_map = build_token_map(summary)

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for token, value in token_map.items():
                        if token in text:
                            text = text.replace(token, value)
                    run.text = text

    prs.save(output_path)
    return output_path

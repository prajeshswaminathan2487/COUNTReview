"""Extract plain text content from an uploaded PPTX or PDF file."""

from pptx import Presentation
from pypdf import PdfReader


def extract_pptx(filepath):
    prs = Presentation(filepath)
    slides_text = []

    for i, slide in enumerate(prs.slides, start=1):
        parts = [f"--- Slide {i} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes: {notes}]")

        slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text)


def extract_pdf(filepath):
    reader = PdfReader(filepath)
    pages_text = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages_text.append(f"--- Page {i} ---\n{text.strip()}")
    return "\n\n".join(pages_text)


def extract_text(filepath, filename):
    lower = filename.lower()
    if lower.endswith(".pptx"):
        return extract_pptx(filepath)
    elif lower.endswith(".pdf"):
        return extract_pdf(filepath)
    else:
        raise ValueError("Unsupported file type. Please upload a .pptx or .pdf file.")
